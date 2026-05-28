import matplotlib.pyplot as plt
from pptx import Presentation
from pptx.util import Inches

# Dados simulados de tendências de mercado
profissoes = [
    "IA/ML Engineer", "Data Scientist", "Cybersecurity Specialist",
    "Cloud Architect", "Software Developer", "DevOps Engineer",
    "UX/UI Designer", "Database Admin", "IT Project Manager", "Infrastructure Analyst"
]

demanda_mercado = [95, 90, 88, 85, 80, 78, 70, 60, 55, 50]  # Valores fictícios de demanda (escala 0–100)

# Gerar o gráfico
plt.figure(figsize=(10, 6))
plt.barh(profissoes, demanda_mercado, color='mediumpurple')
plt.xlabel("Demanda estimada no mercado (%)")
plt.title("Tendência de Profissões em Tecnologia da Informação")
plt.gca().invert_yaxis()  # Inverter para mostrar a profissão mais alta no topo

# Salvar o gráfico
grafico_path = "grafico_profissoes_mercado.png"
plt.tight_layout()
plt.savefig(grafico_path)
plt.close()

# Criar apresentação PowerPoint
prs = Presentation()
slide_layout = prs.slide_layouts[5]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Profissões em Alta no Mercado de TI"

# Inserir gráfico
left = Inches(1)
top = Inches(1.5)
height = Inches(4.5)
slide.shapes.add_picture(grafico_path, left, top, height=height)

# Salvar apresentação
prs.save("Tendencia_Mercado_TI.pptx")
print("Apresentação gerada com sucesso: Tendencia_Mercado_TI.pptx")
